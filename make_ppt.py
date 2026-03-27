from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── 색상 팔레트 (모던 다크 테마) ─────────────────────────────
BG         = RGBColor(0x0D, 0x11, 0x17)   # 거의 검정
SURFACE    = RGBColor(0x16, 0x1B, 0x27)   # 카드 배경
BORDER     = RGBColor(0x21, 0x2D, 0x3F)   # 테두리
BLUE       = RGBColor(0x38, 0x8B, 0xFF)   # 메인 포인트
CYAN       = RGBColor(0x00, 0xD4, 0xFF)   # 강조
GREEN      = RGBColor(0x00, 0xE5, 0xA3)   # 서브 포인트
ORANGE     = RGBColor(0xFF, 0x8C, 0x42)   # 경고/포인트
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
WHITE_DIM  = RGBColor(0xB0, 0xBE, 0xD1)   # 서브 텍스트
DARK_TEXT  = RGBColor(0x7A, 0x8A, 0x9E)   # 더 어두운 서브

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def rect(slide, color, l=0, t=0, w=None, h=None, line_color=None):
    if w is None: w = prs.slide_width
    if h is None: h = prs.slide_height
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s

def rrect(slide, color, l, t, w, h, line_color=None):
    s = slide.shapes.add_shape(9, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h, size=14, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb

def label(slide, text, l, t, color=BLUE, size=10):
    w, h = Inches(1.6), Inches(0.3)
    s = rrect(slide, color, l, t, w, h)
    tf = s.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.runs[0]
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = WHITE
    return s

def accent_line(slide, t, color=BLUE, l=Inches(0.65), w=Inches(12.0)):
    s = slide.shapes.add_shape(1, l, t, w, Pt(1.5))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()

def section_header(slide, title, subtitle=None):
    """공통 헤더: 다크 배경 + 제목 + 하단 포인트 라인"""
    rect(slide, BG)
    rect(slide, SURFACE, 0, 0, prs.slide_width, Inches(1.1))
    accent_line(slide, Inches(1.1), color=BLUE)
    txt(slide, title, Inches(0.65), Inches(0.18), Inches(10), Inches(0.75),
        size=26, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, Inches(0.65), Inches(0.82), Inches(10), Inches(0.35),
            size=11, color=DARK_TEXT, italic=True)

# ═══════════════════════════════════════════════
# SLIDE 1 — 표지
# ═══════════════════════════════════════════════
sl = blank(prs)
rect(sl, BG)

# 왼쪽 포인트 바
rect(sl, BLUE, 0, 0, Inches(0.06), prs.slide_height)

# 중앙 그리드 라인 (디자인 포인트)
for i in range(5):
    s = sl.shapes.add_shape(1, Inches(8.5), Inches(i*1.5), Pt(1), Inches(1.1))
    s.fill.solid(); s.fill.fore_color.rgb = BORDER; s.line.fill.background()

# 태그라인
txt(sl, "PORTFOLIO  ·  2025 – 2026", Inches(0.65), Inches(1.5), Inches(8), Inches(0.4),
    size=11, color=DARK_TEXT, italic=True)

# 메인 타이틀
txt(sl, "PBX 콜센터", Inches(0.65), Inches(2.0), Inches(11), Inches(1.1),
    size=58, bold=True, color=WHITE)
txt(sl, "플랫폼", Inches(0.65), Inches(3.0), Inches(11), Inches(1.1),
    size=58, bold=True, color=CYAN)

# 부제
txt(sl, "Asterisk 20 기반 멀티 테넌트 콜센터 통합 관리 시스템",
    Inches(0.65), Inches(4.15), Inches(9), Inches(0.5),
    size=16, color=WHITE_DIM)

accent_line(sl, Inches(4.85), color=BLUE, l=Inches(0.65), w=Inches(5))

# 기술 뱃지
badges = [
    ("Next.js 14", BLUE),
    ("FastAPI",    RGBColor(0x00, 0x99, 0x77)),
    ("PostgreSQL", RGBColor(0x33, 0x6F, 0xAF)),
    ("Asterisk 20",RGBColor(0x7C, 0x3A, 0xED)),
    ("Docker",     RGBColor(0x02, 0x90, 0xD4)),
]
x = Inches(0.65)
for name, clr in badges:
    label(sl, name, x, Inches(5.2), color=clr, size=10)
    x += Inches(1.75)

txt(sl, "풀스택 개인 프로젝트",
    Inches(0.65), Inches(6.5), Inches(6), Inches(0.4),
    size=12, color=DARK_TEXT)

# ═══════════════════════════════════════════════
# SLIDE 2 — 프로젝트 개요
# ═══════════════════════════════════════════════
sl = blank(prs)
section_header(sl, "프로젝트 개요", "Project Overview")

# 4개 카드
cards = [
    (BLUE,   "🎯  목적",
     "기업 콜센터 운영에 필요한\n내선·IVR·큐·고객·상담 이력을\n단일 플랫폼에서 통합 관리"),
    (GREEN,  "🏢  멀티 테넌트",
     "회사(테넌트)별 독립 데이터 격리\n시스템 관리자 / 관리자 /\n상담원 3단계 역할 체계"),
    (ORANGE, "📞  실시간 통화",
     "Asterisk ARI WebSocket으로\n통화 이벤트 실시간 수신\n통화 상태·이력 자동 기록"),
    (CYAN,   "🔐  권한 관리",
     "계층형 권한 코드 시스템\n메뉴 접근 / 기능 실행\n세분화 제어"),
]
xs = [Inches(0.5), Inches(3.6), Inches(6.7), Inches(9.8)]
for i, (clr, title, body) in enumerate(cards):
    x = xs[i]
    # 왼쪽 포인트 바
    rect(sl, clr, x, Inches(1.3), Inches(0.04), Inches(4.6))
    # 카드 배경
    rrect(sl, SURFACE, x+Inches(0.04), Inches(1.3), Inches(3.0), Inches(4.6))
    txt(sl, title, x+Inches(0.2), Inches(1.5), Inches(2.8), Inches(0.45),
        size=13, bold=True, color=clr)
    txt(sl, body, x+Inches(0.2), Inches(2.05), Inches(2.8), Inches(2.8),
        size=12, color=WHITE_DIM)

# 하단 통계
stats = [("3", "마이크로서비스"), ("17", "DB 테이블"), ("40+", "REST API"), ("10h", "JWT 세션")]
sx = Inches(1.1)
for val, lbl in stats:
    txt(sl, val,  sx, Inches(6.1), Inches(2.7), Inches(0.65),
        size=32, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
    txt(sl, lbl,  sx, Inches(6.75), Inches(2.7), Inches(0.35),
        size=11, color=DARK_TEXT, align=PP_ALIGN.CENTER)
    if sx < Inches(9):
        rect(sl, BORDER, sx+Inches(2.6), Inches(6.2), Pt(1), Inches(0.8))
    sx += Inches(2.85)

# ═══════════════════════════════════════════════
# SLIDE 3 — 시스템 아키텍처
# ═══════════════════════════════════════════════
sl = blank(prs)
section_header(sl, "시스템 아키텍처", "System Architecture")
rect(sl, BG, 0, Inches(1.1), prs.slide_width, prs.slide_height)

layers = [
    (BLUE,   "Browser",    "Next.js 14  ·  React 19  ·  TypeScript  ·  Zustand  ·  Axios  ·  GoJS"),
    (CYAN,   "API Server", "FastAPI  ·  SQLAlchemy 2 (async)  ·  Alembic  ·  JWT HS256  ·  slowapi"),
    (GREEN,  "ARI Worker", "Python asyncio  ·  websockets  ·  httpx  ·  실시간 이벤트 처리"),
    (ORANGE, "Database",   "PostgreSQL 16  ·  pbx_common 공유 ORM 모델  ·  17개 테이블"),
    (RGBColor(0x7C,0x3A,0xED), "Asterisk 20", "PJSIP  ·  ARI (REST Interface)  ·  res_config_pgsql  ·  Docker"),
]

for i, (clr, name, detail) in enumerate(layers):
    t = Inches(1.35) + Inches(1.15) * i
    # 레이어 박스
    rrect(sl, SURFACE, Inches(1.5), t, Inches(10.4), Inches(0.85), line_color=BORDER)
    # 왼쪽 컬러 바
    rect(sl, clr, Inches(1.5), t, Inches(0.07), Inches(0.85))
    # 이름 태그
    txt(sl, name, Inches(1.75), t+Inches(0.08), Inches(1.8), Inches(0.4),
        size=12, bold=True, color=clr)
    # 상세
    txt(sl, detail, Inches(3.6), t+Inches(0.22), Inches(8.0), Inches(0.45),
        size=11, color=WHITE_DIM)
    # 화살표 (마지막 제외)
    if i < 4:
        txt(sl, "↕", Inches(2.05), t+Inches(0.9), Inches(0.4), Inches(0.3),
            size=12, color=BORDER, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════
# SLIDE 4 — 기술 스택
# ═══════════════════════════════════════════════
sl = blank(prs)
section_header(sl, "기술 스택", "Tech Stack")
rect(sl, BG, 0, Inches(1.1), prs.slide_width, prs.slide_height)

cols = [
    (BLUE,   "Frontend", [
        ("Next.js 14",    "App Router · SSR · TypeScript 5"),
        ("React 19",      "컴포넌트 기반 UI"),
        ("Zustand 5",     "전역 상태 (권한·활동상태)"),
        ("Axios 1.13",    "인터셉터로 403 시 권한 자동 갱신"),
        ("GoJS",          "IVR 트리 다이어그램 시각화"),
        ("Tailwind CSS 4","스타일링"),
    ]),
    (GREEN,  "Backend", [
        ("FastAPI",       "Python async REST API"),
        ("SQLAlchemy 2",  "async ORM"),
        ("Alembic",       "DB 마이그레이션"),
        ("JWT HS256",     "10시간 토큰 인증"),
        ("Bcrypt+Pepper", "SHA-256 비밀번호 해싱"),
        ("slowapi",       "Rate Limiting (5/min 로그인)"),
    ]),
    (ORANGE, "Infra / PBX", [
        ("PostgreSQL 16", "관계형 데이터베이스"),
        ("Asterisk 20",   "LTS PBX 엔진"),
        ("Docker Compose","컨테이너 오케스트레이션"),
        ("PJSIP",         "SIP 내선 엔드포인트"),
        ("ARI",           "Asterisk REST Interface"),
        ("res_config_pgsql","DB 기반 설정"),
    ]),
]

xs = [Inches(0.45), Inches(4.65), Inches(8.85)]
for i, (clr, sec, items) in enumerate(cols):
    x = xs[i]
    # 헤더
    rrect(sl, clr, x, Inches(1.25), Inches(3.95), Inches(0.42))
    txt(sl, sec, x, Inches(1.25), Inches(3.95), Inches(0.42),
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # 아이템
    for j, (k, v) in enumerate(items):
        ty = Inches(1.75) + Inches(0.88) * j
        rrect(sl, SURFACE, x, ty, Inches(3.95), Inches(0.78), line_color=BORDER)
        txt(sl, k, x+Inches(0.15), ty+Inches(0.06), Inches(3.65), Inches(0.32),
            size=12, bold=True, color=clr)
        txt(sl, v, x+Inches(0.15), ty+Inches(0.38), Inches(3.65), Inches(0.32),
            size=10.5, color=WHITE_DIM)

# ═══════════════════════════════════════════════
# SLIDE 5 — 주요 기능: IVR & 큐
# ═══════════════════════════════════════════════
sl = blank(prs)
section_header(sl, "주요 기능  ①  —  IVR & 큐 관리", "IVR Tree & Call Queue")
rect(sl, BG, 0, Inches(1.1), prs.slide_width, prs.slide_height)

# IVR 파트
txt(sl, "IVR 흐름 편집기", Inches(0.65), Inches(1.25), Inches(6), Inches(0.4),
    size=14, bold=True, color=BLUE)
accent_line(sl, Inches(1.7), color=BLUE, l=Inches(0.65), w=Inches(5.8))
ivr = [
    "DTMF 키 분기 트리 구조 — parent_id 자기 참조",
    "노드 타입: menu / action / queue 연결",
    "노드별 음성 파일 1:1 업로드",
    "프리셋 플로우 — 전 회사 공용 템플릿",
    "GoJS 캔버스 드래그앤드롭 시각 편집",
    "CASCADE 삭제로 자식 노드 자동 정리",
]
for j, t in enumerate(ivr):
    rrect(sl, SURFACE, Inches(0.65), Inches(1.82)+Inches(0.82)*j,
          Inches(5.8), Inches(0.68), line_color=BORDER)
    rect(sl, BLUE, Inches(0.65), Inches(1.82)+Inches(0.82)*j, Inches(0.04), Inches(0.68))
    txt(sl, t, Inches(0.85), Inches(1.95)+Inches(0.82)*j, Inches(5.5), Inches(0.45),
        size=12, color=WHITE_DIM)

# 큐 파트
txt(sl, "콜 큐 관리", Inches(7.1), Inches(1.25), Inches(5.6), Inches(0.4),
    size=14, bold=True, color=GREEN)
accent_line(sl, Inches(1.7), color=GREEN, l=Inches(7.1), w=Inches(5.6))
queues = [
    "전략: rrmemory (Round-Robin) 등",
    "Ring 타임아웃 / Wrapup 시간 설정",
    "큐 멤버 추가 · 제거 · 일시정지 관리",
    "PJSIP/{exten} 인터페이스 자동 매핑",
    "IVR 노드 → 큐 직접 연결 지원",
]
for j, t in enumerate(queues):
    rrect(sl, SURFACE, Inches(7.1), Inches(1.82)+Inches(0.98)*j,
          Inches(5.6), Inches(0.84), line_color=BORDER)
    rect(sl, GREEN, Inches(7.1), Inches(1.82)+Inches(0.98)*j, Inches(0.04), Inches(0.84))
    txt(sl, t, Inches(7.3), Inches(2.02)+Inches(0.98)*j, Inches(5.3), Inches(0.45),
        size=12, color=WHITE_DIM)

# ═══════════════════════════════════════════════
# SLIDE 6 — 주요 기능: 상담 처리
# ═══════════════════════════════════════════════
sl = blank(prs)
section_header(sl, "주요 기능  ②  —  실시간 상담 처리", "Real-time Consultation")
rect(sl, BG, 0, Inches(1.1), prs.slide_width, prs.slide_height)

# 상담 흐름 스텝
steps = [
    (BLUE,   "01", "콜 인입",       "ARI WebSocket\n대기열 수신"),
    (CYAN,   "02", "고객 조회",     "전화번호 → DB 검색\n프로필 자동 표시"),
    (GREEN,  "03", "고객 등록",     "미등록 시 즉시 등록\n이름/전화/그룹"),
    (ORANGE, "04", "메모 작성",     "카테고리 3단계 선택\n자유 메모 입력"),
    (RGBColor(0xEC,0x4E,0x99), "05", "저장·종료", "상담 이력 DB 저장\n타이머 자동 기록"),
]
sx = Inches(0.5)
for clr, num, title, desc in steps:
    rrect(sl, SURFACE, sx, Inches(1.3), Inches(2.35), Inches(2.0), line_color=BORDER)
    rect(sl, clr, sx, Inches(1.3), prs.slide_width, Inches(0.04))  # 상단 컬러 라인
    rrect(sl, clr, sx+Inches(0.15), Inches(1.3), Inches(2.35), Inches(0.04))
    txt(sl, num, sx+Inches(0.15), Inches(1.38), Inches(2.05), Inches(0.45),
        size=22, bold=True, color=clr)
    txt(sl, title, sx+Inches(0.15), Inches(1.82), Inches(2.05), Inches(0.4),
        size=13, bold=True, color=WHITE)
    txt(sl, desc, sx+Inches(0.15), Inches(2.28), Inches(2.05), Inches(0.85),
        size=11, color=WHITE_DIM)
    if sx < Inches(10):
        txt(sl, "→", sx+Inches(2.35), Inches(2.0), Inches(0.3), Inches(0.4),
            size=16, color=BORDER, align=PP_ALIGN.CENTER)
    sx += Inches(2.6)

# 버전 관리 섹션
accent_line(sl, Inches(3.65), color=BLUE, l=Inches(0.65), w=Inches(12.0))
txt(sl, "상담 이력 버전 관리 (Immutable History)", Inches(0.65), Inches(3.75),
    Inches(8), Inches(0.4), size=14, bold=True, color=BLUE)

ver_items = [
    ("수정", ORANGE, "기존 레코드 status → SUPERSEDED  +  신규 ACTIVE 레코드 생성 (original_id 참조)"),
    ("삭제", RGBColor(0xEC,0x4E,0x99), "status → INACTIVE  (물리 삭제 없음, 이력 완전 보존)"),
    ("조회", GREEN, "기본값: INACTIVE 제외  /  status 파라미터로 특정 버전 직접 조회 가능"),
]
for j, (tag, clr, detail) in enumerate(ver_items):
    ty = Inches(4.3) + Inches(0.88) * j
    rrect(sl, SURFACE, Inches(0.65), ty, Inches(12.0), Inches(0.72), line_color=BORDER)
    rrect(sl, clr, Inches(0.65), ty, Inches(0.9), Inches(0.72))
    txt(sl, tag, Inches(0.65), ty, Inches(0.9), Inches(0.72),
        size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, detail, Inches(1.7), ty+Inches(0.15), Inches(10.8), Inches(0.45),
        size=11.5, color=WHITE_DIM)

# ═══════════════════════════════════════════════
# SLIDE 7 — 보안 & 인증
# ═══════════════════════════════════════════════
sl = blank(prs)
section_header(sl, "보안 & 인증 아키텍처", "Security & Auth")
rect(sl, BG, 0, Inches(1.1), prs.slide_width, prs.slide_height)

# 비밀번호 파이프라인
txt(sl, "비밀번호 해싱 파이프라인", Inches(0.65), Inches(1.3), Inches(8), Inches(0.4),
    size=13, bold=True, color=BLUE)
pipeline = ["평문 비밀번호", "Pepper 결합", "SHA-256", "Base64 인코딩", "Bcrypt 해싱"]
px = Inches(0.65)
for i, step in enumerate(pipeline):
    clr = BLUE if i < 4 else GREEN
    rrect(sl, clr, px, Inches(1.82), Inches(2.1), Inches(0.5))
    txt(sl, step, px, Inches(1.82), Inches(2.1), Inches(0.5),
        size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < 4:
        txt(sl, "→", px+Inches(2.12), Inches(1.88), Inches(0.28), Inches(0.35),
            size=16, color=CYAN, align=PP_ALIGN.CENTER)
    px += Inches(2.42)

# JWT 페이로드
txt(sl, "JWT 페이로드", Inches(0.65), Inches(2.65), Inches(5.5), Inches(0.4),
    size=13, bold=True, color=BLUE)
jwt_fields = [
    ("sub", "계정 (account)"),
    ("name", "사용자 이름"),
    ("role", "SYSTEM_ADMIN / MANAGER / AGENT"),
    ("id", "user.id"),
    ("company_id", "소속 회사 ID"),
    ("exp", "발급 후 10시간"),
]
for j, (k, v) in enumerate(jwt_fields):
    ty = Inches(3.15) + Inches(0.63) * j
    rrect(sl, SURFACE, Inches(0.65), ty, Inches(5.5), Inches(0.55), line_color=BORDER)
    txt(sl, k, Inches(0.8), ty+Inches(0.08), Inches(1.4), Inches(0.38),
        size=11.5, bold=True, color=CYAN)
    txt(sl, v, Inches(2.35), ty+Inches(0.08), Inches(3.65), Inches(0.38),
        size=11, color=WHITE_DIM)

# 권한 처리 흐름
txt(sl, "권한 처리 흐름", Inches(7.1), Inches(2.65), Inches(5.7), Inches(0.4),
    size=13, bold=True, color=GREEN)
auth_steps = [
    "① 로그인 시 권한 코드 목록 수신",
    "② Zustand authStore 저장 (10h TTL)",
    "③ hasPermission('code') 컴포넌트 단 제어",
    "④ API 요청 시 Bearer 토큰 헤더 첨부",
    "⑤ 백엔드 require_permission() 데코레이터 검증",
    "⑥ 403 응답 시 /auth/me/permissions 자동 재조회",
]
for j, t in enumerate(auth_steps):
    ty = Inches(3.15) + Inches(0.63) * j
    rrect(sl, SURFACE, Inches(7.1), ty, Inches(5.65), Inches(0.55), line_color=BORDER)
    txt(sl, t, Inches(7.25), ty+Inches(0.1), Inches(5.4), Inches(0.38),
        size=11.5, color=WHITE_DIM)

# ═══════════════════════════════════════════════
# SLIDE 8 — DB 스키마
# ═══════════════════════════════════════════════
sl = blank(prs)
section_header(sl, "데이터베이스 스키마", "Database Schema  ·  17 Tables")
rect(sl, BG, 0, Inches(1.1), prs.slide_width, prs.slide_height)

tables = [
    (BLUE,   "company",          "id · company_name · manager_name · use_callback · is_active · business_number"),
    (BLUE,   "user",             "id · company_id · account · account_pw · exten · name · role · is_active"),
    (GREEN,  "customer",         "id · company_id · name · phone · group(vip/normal/blacklist) · is_active · last_call_at"),
    (ORANGE, "calls",            "id(UUID) · caller_exten · callee_exten · started_at · answered_at · ended_at · direction"),
    (ORANGE, "call_events",      "id · call_id · ts · type · channel_id · raw(JSONB)"),
    (RGBColor(0x7C,0x3A,0xED), "queues", "id · company_id · name · strategy · timeout · wrapuptime · maxlen"),
    (RGBColor(0x7C,0x3A,0xED), "ivr_flows",  "id · company_id · name · is_preset · is_active"),
    (RGBColor(0x7C,0x3A,0xED), "ivr_nodes",  "id · flow_id · parent_id(자기참조) · queue_id · dtmf_key · node_type · config(JSONB)"),
    (CYAN,   "consultations",    "id · call_id · agent_id · original_id · category_id · memo · status(ACTIVE/SUPERSEDED/INACTIVE)"),
    (CYAN,   "consult_categories","id · company_id · parent_id(3단계) · name · depth"),
    (DARK_TEXT,"permissions",    "id · parent_id · code(unique) · name · type(MENU/ACTION)"),
    (DARK_TEXT,"user_status",    "user_id(PK/FK) · login_status · activity · last_login_at"),
]

for j, (clr, tname, fields) in enumerate(tables):
    ty = Inches(1.3) + Inches(0.5) * j
    rrect(sl, SURFACE, Inches(0.5), ty, Inches(12.4), Inches(0.42), line_color=BORDER)
    rrect(sl, clr, Inches(0.5), ty, Inches(1.9), Inches(0.42))
    txt(sl, tname, Inches(0.5), ty, Inches(1.9), Inches(0.42),
        size=10.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, fields, Inches(2.5), ty+Inches(0.05), Inches(10.2), Inches(0.33),
        size=10, color=WHITE_DIM)

# ═══════════════════════════════════════════════
# SLIDE 9 — ARI Worker
# ═══════════════════════════════════════════════
sl = blank(prs)
section_header(sl, "Asterisk ARI Worker", "실시간 통화 이벤트 처리")
rect(sl, BG, 0, Inches(1.1), prs.slide_width, prs.slide_height)

# 파이프라인
pipeline_nodes = [
    (BG_DARK if False else SURFACE, CYAN,   "Asterisk 20\nARI WebSocket"),
    (SURFACE, BLUE,   "parser.py\nParseEvent 파싱"),
    (SURFACE, GREEN,  "call_service.py\n이벤트 처리 & 상태 추적"),
    (SURFACE, ORANGE, "call_recorder.py\nPostgreSQL 저장"),
]
# 배경 연결 선
rect(sl, BORDER, Inches(0.65), Inches(2.5), Inches(12.0), Pt(2))

nx = Inches(0.65)
for i, (bg, clr, label_txt) in enumerate(pipeline_nodes):
    rrect(sl, SURFACE, nx, Inches(1.85), Inches(2.6), Inches(1.3), line_color=clr)
    rect(sl, clr, nx, Inches(1.85), Inches(2.6), Inches(0.06))
    txt(sl, label_txt, nx+Inches(0.1), Inches(2.1), Inches(2.4), Inches(0.9),
        size=12, bold=True, color=clr, align=PP_ALIGN.CENTER)
    if i < 3:
        txt(sl, "→", nx+Inches(2.62), Inches(2.3), Inches(0.38), Inches(0.4),
            size=18, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
    nx += Inches(3.15)

# 이벤트 목록
txt(sl, "처리 ARI 이벤트", Inches(0.65), Inches(3.45), Inches(5), Inches(0.4),
    size=13, bold=True, color=BLUE)
events = [
    ("StasisStart",        BLUE,   "신규 채널 진입 — 통화 추적 시작"),
    ("ChannelStateChange", CYAN,   "채널 상태 변경 — 통화 상태 업데이트"),
    ("BridgeCreated",      GREEN,  "브릿지 생성 — 양방향 통화 확립"),
    ("ChannelLeftBridge",  ORANGE, "브릿지 이탈 — 통화 종료 처리"),
    ("ChannelHangup",      RGBColor(0xEC,0x4E,0x99), "채널 Hangup — 종료 사유 기록"),
]
for j, (ev, clr, desc) in enumerate(events):
    ty = Inches(3.98) + Inches(0.65) * j
    rrect(sl, SURFACE, Inches(0.65), ty, Inches(12.0), Inches(0.55), line_color=BORDER)
    rrect(sl, clr, Inches(0.65), ty, Inches(2.4), Inches(0.55))
    txt(sl, ev, Inches(0.65), ty, Inches(2.4), Inches(0.55),
        size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, desc, Inches(3.2), ty+Inches(0.1), Inches(9.2), Inches(0.38),
        size=11.5, color=WHITE_DIM)

# ═══════════════════════════════════════════════
# SLIDE 10 — 마무리
# ═══════════════════════════════════════════════
sl = blank(prs)
rect(sl, BG)
rect(sl, BLUE, 0, 0, Inches(0.06), prs.slide_height)

txt(sl, "풀스택 단독 설계 · 구현", Inches(0.65), Inches(1.6), Inches(12), Inches(0.8),
    size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(sl, "Frontend  ·  Backend  ·  ARI Worker  ·  Infra",
    Inches(0.65), Inches(2.5), Inches(12), Inches(0.5),
    size=20, color=CYAN, align=PP_ALIGN.CENTER)

accent_line(sl, Inches(3.3), color=BORDER, l=Inches(2.5), w=Inches(8.3))

highlights = [
    (BLUE,   "Next.js 14 App Router  ·  FastAPI async  ·  Asterisk 20 ARI 실시간 연동"),
    (GREEN,  "멀티 테넌트 접근 제어  ·  계층형 권한 시스템  ·  상담 이력 버전 관리"),
    (ORANGE, "Docker 기반 인프라  ·  JWT 인증  ·  Bcrypt+Pepper+SHA-256 보안"),
]
for j, (clr, t) in enumerate(highlights):
    ty = Inches(3.6) + Inches(0.88) * j
    rrect(sl, SURFACE, Inches(1.5), ty, Inches(10.3), Inches(0.68), line_color=BORDER)
    rect(sl, clr, Inches(1.5), ty, Inches(0.05), Inches(0.68))
    txt(sl, t, Inches(1.7), ty+Inches(0.14), Inches(10.0), Inches(0.42),
        size=13, color=WHITE_DIM, align=PP_ALIGN.CENTER)

txt(sl, "github.com/jang-yeong-wan",
    Inches(0.65), Inches(6.7), Inches(12), Inches(0.4),
    size=12, color=DARK_TEXT, italic=True, align=PP_ALIGN.CENTER)

# ── 저장 ──────────────────────────────────────────────────────
output = "/Users/jang-yeong-wan/PBX_asterisk/PBX_Platform_Portfolio.pptx"
prs.save(output)
print(f"✅  저장 완료: {output}")
